# src/ingestion/pptx_parser.py

from pptx import Presentation
from ingestion_utils import save_text

def parse_pptx(file_path: str, output_dir: str, document_id: str) -> str:
    """
    Extract text from PPTX slides and save it.
    """
    prs = Presentation(file_path)
    full_text = ""
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"
    
    output_path = save_text(full_text, output_dir, document_id)
    return output_path
